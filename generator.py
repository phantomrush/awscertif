from pptx import Presentation

search_str1 = 'pgpplname'
repl_str1 = 'Saksham Singal'
search_str2 = 'pgname'
repl_str2 = 'FLY Find the Leader in You'
search_str3 = 'pgdate'
repl_str3 = '12th Aug, 2018'
ppt = Presentation('tgen.pptx')

for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            paragraphs = shape.text_frame.paragraphs
            print(paragraphs)
            for paras in paragraphs:
                for oldr in paras.runs:
                    print(oldr.text)
                    oldr.text = oldr.text.replace(search_str1, repl_str1)
                    oldr.text = oldr.text.replace(search_str2, repl_str2)
                    oldr.text = oldr.text.replace(search_str3, repl_str3)
                    print(oldr.text)
ppt.save('certificate.pptx')
